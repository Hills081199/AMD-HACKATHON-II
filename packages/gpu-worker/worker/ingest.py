"""Pipeline step 1 — chunk documents by heading/slide boundary.

See docs/concept-graph-pipeline.md step 1. Deliberately does not chunk by a
hard character count: a fixed-size window can cut a concept in half, which
poisons concept extraction (step 2) downstream.

IMP-B1 improvements:
- PDF: merge short pages (<200 chars) with previous; split long pages (>2000
  chars) at double-newline boundaries; adds 1-sentence overlap between chunks.
- DOCX: splits on Heading 1/2/3 (not just H1), so subsections get their own
  chunk.
- PPTX: unchanged (slide boundary is already the right granularity).
"""

from __future__ import annotations

import io
import re
from dataclasses import asdict, dataclass
from pathlib import PurePosixPath

_MIN_CHUNK_CHARS = 200   # merge shorter chunks into the previous one
_MAX_CHUNK_CHARS = 2000  # split longer chunks at paragraph boundary
_OVERLAP_SENTENCES = 1   # sentences to carry over between PDF page chunks


@dataclass
class Chunk:
    doc_id: str
    chunk_id: str
    page: int
    text: str

    def to_dict(self) -> dict:
        return asdict(self)


def chunk_document(filename: str, content: bytes) -> list[Chunk]:
    """Dispatch to the right parser based on file extension and return
    boundary-aligned chunks for one document."""
    doc_id = filename
    suffix = PurePosixPath(filename).suffix.lower()
    if suffix == ".pdf":
        return _chunk_pdf(doc_id, content)
    if suffix == ".pptx":
        return _chunk_pptx(doc_id, content)
    if suffix == ".docx":
        return _chunk_docx(doc_id, content)
    raise ValueError(f"Unsupported document type: {suffix or filename!r}")


def _split_long_text(text: str, max_chars: int = _MAX_CHUNK_CHARS) -> list[str]:
    """Split text that exceeds max_chars at double-newline paragraph boundaries.
    Falls back to single-newline, then sentence boundary if no paragraph break
    exists within the limit."""
    if len(text) <= max_chars:
        return [text]

    parts: list[str] = []
    remaining = text
    while len(remaining) > max_chars:
        # Try to find a paragraph break before the limit
        split_pos = remaining.rfind("\n\n", 0, max_chars)
        if split_pos == -1:
            # Fall back to single newline
            split_pos = remaining.rfind("\n", 0, max_chars)
        if split_pos == -1:
            # Last resort: split at sentence boundary (period + space)
            split_pos = remaining.rfind(". ", 0, max_chars)
            if split_pos != -1:
                split_pos += 1  # include the period
        if split_pos <= 0:
            # No good split point — hard cut at limit
            split_pos = max_chars
        parts.append(remaining[:split_pos].strip())
        remaining = remaining[split_pos:].strip()
    if remaining:
        parts.append(remaining)
    return [p for p in parts if p]


def _last_sentence(text: str) -> str:
    """Return the last complete sentence from text for overlap."""
    # Split on period/exclamation/question followed by whitespace or end
    sentences = re.split(r'(?<=[.!?])\s+', text.strip())
    if not sentences:
        return ""
    return sentences[-1].strip()


def _chunk_pdf(doc_id: str, content: bytes) -> list[Chunk]:
    import fitz  # PyMuPDF

    raw_pages: list[tuple[int, str]] = []
    with fitz.open(stream=content, filetype="pdf") as pdf:
        for page_number, page in enumerate(pdf, start=1):
            text = page.get_text().strip()
            if text:
                raw_pages.append((page_number, text))

    if not raw_pages:
        return []

    # Merge short pages into the previous page's text
    merged: list[tuple[int, str]] = []
    for page_number, text in raw_pages:
        if merged and len(text) < _MIN_CHUNK_CHARS:
            prev_num, prev_text = merged[-1]
            merged[-1] = (prev_num, prev_text + "\n\n" + text)
        else:
            merged.append((page_number, text))

    chunks: list[Chunk] = []
    prev_overlap = ""
    chunk_counter = 0

    for page_number, text in merged:
        # Prepend overlap from previous chunk
        full_text = (prev_overlap + "\n" + text).strip() if prev_overlap else text

        # Split long pages
        parts = _split_long_text(full_text)
        for part in parts:
            if not part:
                continue
            chunk_counter += 1
            chunks.append(
                Chunk(
                    doc_id=doc_id,
                    chunk_id=f"{doc_id}:p{chunk_counter}",
                    page=page_number,
                    text=part,
                )
            )
        # Carry overlap from the last part of this page
        prev_overlap = _last_sentence(parts[-1]) if parts else ""

    return chunks


def _chunk_pptx(doc_id: str, content: bytes) -> list[Chunk]:
    from pptx import Presentation

    chunks: list[Chunk] = []
    presentation = Presentation(io.BytesIO(content))
    for slide_number, slide in enumerate(presentation.slides, start=1):
        parts = [
            shape.text_frame.text.strip()
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        text = "\n".join(parts)
        if not text:
            continue
        chunks.append(
            Chunk(doc_id=doc_id, chunk_id=f"{doc_id}:s{slide_number}", page=slide_number, text=text)
        )
    return chunks


def _chunk_docx(doc_id: str, content: bytes) -> list[Chunk]:
    """Split DOCX on Heading 1/2/3 (IMP-B1: was Heading 1 only), so subsections
    get their own chunk instead of being merged into a giant parent section.
    `page` here is the 1-indexed heading-delimited section, which is the only
    stable locator available without a rendering pass."""
    from docx import Document

    document = Document(io.BytesIO(content))
    chunks: list[Chunk] = []
    section_number = 0
    buffer: list[str] = []

    def flush() -> None:
        nonlocal section_number
        text = "\n".join(buffer).strip()
        if text and len(text) >= _MIN_CHUNK_CHARS:
            section_number += 1
            # Split oversized sections
            parts = _split_long_text(text)
            for part in parts:
                chunks.append(
                    Chunk(
                        doc_id=doc_id,
                        chunk_id=f"{doc_id}:h{section_number}",
                        page=section_number,
                        text=part,
                    )
                )
        elif text:
            # Too short — append to last chunk if possible, otherwise include anyway
            if chunks:
                last = chunks[-1]
                chunks[-1] = Chunk(
                    doc_id=last.doc_id,
                    chunk_id=last.chunk_id,
                    page=last.page,
                    text=last.text + "\n" + text,
                )
            else:
                section_number += 1
                chunks.append(
                    Chunk(doc_id=doc_id, chunk_id=f"{doc_id}:h{section_number}", page=section_number, text=text)
                )
        buffer.clear()

    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        style_name = paragraph.style.name.lower()
        # IMP-B1: split on H1, H2, H3 instead of just H1
        is_heading = style_name.startswith("heading") and any(
            style_name.startswith(f"heading {n}") or style_name == f"heading{n}"
            for n in ("1", "2", "3", " 1", " 2", " 3")
        ) or style_name.startswith("title")
        if is_heading and buffer:
            flush()
        buffer.append(text)
    flush()
    return chunks
