"""Tests for pipeline step 1 (docs/concept-graph-pipeline.md).

Builds synthetic PDF/PPTX/DOCX fixtures in-memory since no sample corpus is
checked in yet (data/ only has the mastery-tree sample JSON) — these fixtures
are enough to prove the boundary-splitting behavior the feature depends on.
"""

import io

import fitz  # PyMuPDF
import pytest
from docx import Document
from pptx import Presentation

from worker.ingest import chunk_document


def _make_pdf_bytes(pages: list[str]) -> bytes:
    doc = fitz.open()
    for text in pages:
        page = doc.new_page()
        page.insert_text((72, 72), text)
    data = doc.tobytes()
    doc.close()
    return data


def _make_pptx_bytes(slides: list[str]) -> bytes:
    presentation = Presentation()
    blank_layout = presentation.slide_layouts[6]
    for text in slides:
        slide = presentation.slides.add_slide(blank_layout)
        textbox = slide.shapes.add_textbox(0, 0, presentation.slide_width, presentation.slide_height)
        textbox.text_frame.text = text
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _make_docx_bytes(sections: list[tuple[str, str]]) -> bytes:
    document = Document()
    for heading, body in sections:
        document.add_heading(heading, level=1)
        document.add_paragraph(body)
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def test_chunk_pdf_splits_by_page():
    content = _make_pdf_bytes(["Linear algebra basics.", "Gradient descent intro."])
    chunks = chunk_document("linear_algebra.pdf", content)

    assert [c.page for c in chunks] == [1, 2]
    assert all(c.doc_id == "linear_algebra.pdf" for c in chunks)
    assert len({c.chunk_id for c in chunks}) == 2
    assert "Linear algebra" in chunks[0].text
    assert "Gradient descent" in chunks[1].text


def test_chunk_pptx_splits_by_slide():
    content = _make_pptx_bytes(["Derivatives", "Partial derivatives"])
    chunks = chunk_document("calculus.pptx", content)

    assert [c.page for c in chunks] == [1, 2]
    assert chunks[0].text == "Derivatives"
    assert chunks[1].text == "Partial derivatives"


def test_chunk_docx_splits_by_heading_not_mid_section():
    content = _make_docx_bytes(
        [
            ("Overfitting", "A model that fits noise instead of signal."),
            ("Regularization", "L1 and L2 penalties reduce overfitting."),
        ]
    )
    chunks = chunk_document("ml_notes.docx", content)

    assert len(chunks) == 2
    assert chunks[0].text.startswith("Overfitting")
    assert "L1 and L2" not in chunks[0].text
    assert chunks[1].text.startswith("Regularization")


def test_unsupported_extension_raises():
    with pytest.raises(ValueError):
        chunk_document("notes.txt", b"plain text")
