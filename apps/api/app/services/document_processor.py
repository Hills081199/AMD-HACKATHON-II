"""
Document processing service — Pipeline Steps 1-6.

Replaces the old single-LLM-call approach with the full gpu-worker pipeline:
  Step 1: chunk_document()          [worker/ingest.py]
  Step 2: extract_raw_concepts()    [worker/concepts.py]
  Step 3: cluster_concepts()        [worker/concepts.py]
  Step 4: build_candidate_edges()   [worker/prerequisites.py]
  Step 5: validate_graph()          [services/graph.py]
  Step 6: assign_levels()           [services/levels.py]

Step 7 (lesson / quiz generation) is intentionally omitted here — each node
will get its own API endpoint for on-demand generation later.
"""
from __future__ import annotations

import os
import sys
import re
import json
from datetime import datetime
from pathlib import Path
from typing import Any, Callable

# ---------------------------------------------------------------------------
# Bootstrap: make gpu-worker importable from inside apps/api.
# Works in 3 environments:
#   1. Local dev  : repo/apps/api/app/services/ → parents[4] = repo root
#   2. Docker     : /app/app/services/          → GPU_WORKER_PATH env var
#                   or /gpu-worker/ volume mount
#   3. Any other  : GPU_WORKER_PATH env var
# ---------------------------------------------------------------------------
def _find_gpu_worker() -> Path | None:
    # Strategy 1: explicit env var (most reliable in Docker)
    env_path = os.getenv("GPU_WORKER_PATH")
    if env_path:
        p = Path(env_path)
        if p.exists():
            return p

    # Strategy 2: walk up from this file looking for packages/gpu-worker
    here = Path(__file__).resolve()
    for ancestor in here.parents:
        candidate = ancestor / "packages" / "gpu-worker"
        if candidate.exists():
            return candidate

    # Strategy 3: known Docker volume-mount location
    docker_path = Path("/gpu-worker")
    if docker_path.exists():
        return docker_path

    return None


_GPU_WORKER = _find_gpu_worker()
if _GPU_WORKER is not None and str(_GPU_WORKER) not in sys.path:
    sys.path.insert(0, str(_GPU_WORKER))

from worker.ingest import Chunk, chunk_document  # noqa: E402
from worker.concepts import (  # noqa: E402
    OpenAIConceptExtractor,
    GemmaConceptExtractor,
    OpenAIEmbedder,
    SentenceTransformerEmbedder,
    cluster_concepts,
    extract_raw_concepts,
)
from worker.prerequisites import FireworksClient as PrereqFireworksClient, build_candidate_edges  # noqa: E402

from app.services.graph import validate_graph  # noqa: E402
from app.services.levels import assign_levels  # noqa: E402

# ---------------------------------------------------------------------------
# Config
# ---------------------------------------------------------------------------
LLM_PROVIDER = os.getenv("LLM_PROVIDER", "openai").lower()
PREREQ_SIMILARITY_THRESHOLD = float(os.getenv("PREREQ_SIMILARITY_THRESHOLD", "0.6"))
# P99.5 of embedding distribution → top 0.5% most-similar pairs treated as near-duplicates
# Matches build_demo_dataset.py default; decrease to merge more aggressively
DEDUP_PERCENTILE = float(os.getenv("DEDUP_PERCENTILE", "99.5"))
MIN_DEPTH_WARN = int(os.getenv("MIN_DEPTH_WARN", "3"))


def _make_extractor():
    if LLM_PROVIDER == "openai":
        return OpenAIConceptExtractor()
    base_url = os.getenv("GEMMA_BASE_URL", "http://localhost:11434")
    return GemmaConceptExtractor(base_url=base_url)


def _make_embedder():
    if LLM_PROVIDER == "openai":
        return OpenAIEmbedder()
    return SentenceTransformerEmbedder()


def _make_prereq_client():
    # FireworksClient auto-switches to OpenAI when LLM_PROVIDER=openai
    return PrereqFireworksClient()


# ---------------------------------------------------------------------------
# Legacy text extraction (kept for backward compat / fallback)
# ---------------------------------------------------------------------------
import fitz  # PyMuPDF
from pptx import Presentation
from docx import Document as DocxDocument


def extract_text_from_pdf(file_path: str) -> str:
    text_parts = []
    try:
        doc = fitz.open(file_path)
        for page_num, page in enumerate(doc):
            text = page.get_text()
            if text.strip():
                text_parts.append(f"[Page {page_num + 1}]\n{text}")
        doc.close()
    except Exception as e:
        raise Exception(f"Failed to extract PDF: {str(e)}")
    return "\n\n".join(text_parts)


def extract_text_from_pptx(file_path: str) -> str:
    text_parts = []
    try:
        prs = Presentation(file_path)
        for slide_num, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            if slide_text:
                text_parts.append(f"[Slide {slide_num + 1}]\n" + "\n".join(slide_text))
    except Exception as e:
        raise Exception(f"Failed to extract PPTX: {str(e)}")
    return "\n\n".join(text_parts)


def extract_text_from_docx(file_path: str) -> str:
    text_parts = []
    try:
        doc = DocxDocument(file_path)
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
    except Exception as e:
        raise Exception(f"Failed to extract DOCX: {str(e)}")
    return "\n\n".join(text_parts)


def extract_text(file_path: str, file_type: str) -> str:
    if file_type == "pdf":
        return extract_text_from_pdf(file_path)
    elif file_type == "pptx":
        return extract_text_from_pptx(file_path)
    elif file_type == "docx":
        return extract_text_from_docx(file_path)
    else:
        raise ValueError(f"Unsupported file type: {file_type}")


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
_DIFFICULTY_BADGE = {
    "foundational": "⭐ Foundational",
    "intermediate": "⭐⭐ Intermediate",
    "advanced": "⭐⭐⭐ Advanced",
}
_LEVEL_XP = {0: 50, 1: 100, 2: 150}
_LEVEL_MINUTES = {0: 10, 1: 15, 2: 20}


def _slugify(name: str) -> str:
    return re.sub(r"[^a-z0-9]+", "_", name.lower()).strip("_")


def _xp_reward(level: int) -> int:
    return _LEVEL_XP.get(level, 150 + (level - 2) * 25)


def _estimated_minutes(level: int) -> int:
    return _LEVEL_MINUTES.get(level, 20 + (level - 2) * 5)


# ---------------------------------------------------------------------------
# Position calculation
# ---------------------------------------------------------------------------
def calculate_node_positions(nodes: list[dict], edges: list[dict]) -> list[dict]:
    levels: dict[int, list] = {}
    for node in nodes:
        level = node.get("level", 0)
        levels.setdefault(level, []).append(node)

    y_spacing = 160
    x_spacing = 220

    for level, level_nodes in levels.items():
        y = level * y_spacing
        total_width = (len(level_nodes) - 1) * x_spacing
        start_x = -total_width / 2
        for i, node in enumerate(level_nodes):
            node["position"] = {"x": start_x + i * x_spacing, "y": y}

    return nodes


# ---------------------------------------------------------------------------
# Main pipeline — Steps 1-6
# ---------------------------------------------------------------------------
async def process_documents(
    topic_id: str,
    file_paths: list[dict[str, str]],
    progress_callback: Callable[[str, str], None] | None = None,
) -> dict[str, Any]:
    """
    Run the full pipeline (Steps 1-6) over the uploaded documents.

    Args:
        topic_id: UUID of the topic (used for metadata only).
        file_paths: List of dicts with keys 'path', 'filename', 'file_type'.
        progress_callback: Optional callable(step_label, detail) called at
                           each pipeline step so callers can update the DB.

    Returns:
        Complete mastery tree dict with keys: topic, nodes, edges, …
    """

    def _emit(step: str, detail: str = "") -> None:
        print(f"[Pipeline] {step}: {detail}")
        if progress_callback:
            progress_callback(step, detail)

    # ------------------------------------------------------------------
    # Step 1 — Chunk documents
    # ------------------------------------------------------------------
    _emit("chunking", f"Processing {len(file_paths)} document(s)…")
    chunks: list[Chunk] = []
    for fi in file_paths:
        try:
            content = Path(fi["path"]).read_bytes()
            new_chunks = chunk_document(fi["filename"], content)
            chunks.extend(new_chunks)
            _emit("chunking", f"  {fi['filename']} → {len(new_chunks)} chunk(s)")
        except Exception as exc:
            print(f"[Pipeline] Warning: failed to chunk {fi['filename']}: {exc}")

    if not chunks:
        raise ValueError("No chunks could be produced from the uploaded documents.")

    _emit("chunking", f"Total {len(chunks)} chunk(s) produced.")

    # ------------------------------------------------------------------
    # Step 2 — Extract raw concepts
    # ------------------------------------------------------------------
    _emit("extracting", f"Extracting concepts from {len(chunks)} chunk(s)…")
    extractor = _make_extractor()
    raw_concepts = extract_raw_concepts(chunks, extractor)
    _emit("extracting", f"{len(raw_concepts)} raw concept(s) found.")

    # ------------------------------------------------------------------
    # Step 3 — Cluster / deduplicate
    # ------------------------------------------------------------------
    _emit("clustering", f"Clustering {len(raw_concepts)} concept(s) (adaptive P{DEDUP_PERCENTILE} threshold)…")
    embedder = _make_embedder()
    canonical_concepts = cluster_concepts(raw_concepts, embedder, dedup_percentile=DEDUP_PERCENTILE)
    concept_dicts = [c.to_dict() for c in canonical_concepts]

    # Difficulty distribution diagnostic (mirrors build_demo_dataset)
    diff_counts: dict[str, int] = {}
    for rc in raw_concepts:
        diff_counts[rc.difficulty] = diff_counts.get(rc.difficulty, 0) + 1
    print(f"[Pipeline] Difficulty distribution: {diff_counts}")

    _emit("clustering", f"Reduced to {len(concept_dicts)} canonical concept(s).")

    # ------------------------------------------------------------------
    # Step 4 — Infer prerequisite edges
    # ------------------------------------------------------------------
    _emit("inferring", f"Inferring prerequisites for {len(concept_dicts)} concept(s)…")
    prereq_client = _make_prereq_client()
    candidate_edges = build_candidate_edges(
        concept_dicts,
        prereq_client,
        similarity_threshold=PREREQ_SIMILARITY_THRESHOLD,
    )
    _emit("inferring", f"{len(candidate_edges)} candidate edge(s) found.")

    # ------------------------------------------------------------------
    # Step 5 — Validate DAG (cycle detection & repair)
    # ------------------------------------------------------------------
    _emit("validating", "Validating dependency graph…")
    validated = validate_graph(candidate_edges)
    _emit(
        "validating",
        f"Dropped {len(validated['dropped_edges'])} edge(s) to fix cycles. "
        f"Valid: {len(validated['edges'])} edge(s).",
    )

    # ------------------------------------------------------------------
    # Step 6 — Assign levels
    # ------------------------------------------------------------------
    _emit("leveling", "Assigning concept tiers…")
    leveled_nodes = assign_levels(validated["edges"], concept_dicts)
    max_level = max((n["level"] for n in leveled_nodes), default=0)

    # IMP-D2: warn if tree is too shallow (mirrors build_demo_dataset)
    if max_level < MIN_DEPTH_WARN:
        print(
            f"[Pipeline] WARNING: Tree depth is only {max_level + 1} level(s) — expected at least {MIN_DEPTH_WARN}.\n"
            f"  Try: lower PREREQ_SIMILARITY_THRESHOLD (current: {PREREQ_SIMILARITY_THRESHOLD}) "
            f"or raise PREREQ_MAX_ALL_PAIRS env var."
        )

    _emit("leveling", f"Tree depth: {max_level + 1} level(s), {len(leveled_nodes)} node(s).")

    # ------------------------------------------------------------------
    # Build output (Step 7 skipped — quiz/lesson generated per-node later)
    # ------------------------------------------------------------------
    _emit("building", "Assembling tree output…")

    # chunk_id → chunk for source lookup
    chunks_by_id = {c.chunk_id: c for c in chunks}

    # Prerequisite IDs per node
    prerequisites_by_id: dict[str, list[str]] = {n["id"]: [] for n in leveled_nodes}
    for edge in validated["edges"]:
        prerequisites_by_id.setdefault(edge["to"], []).append(edge["from"])

    nodes_out: list[dict] = []
    for leveled in leveled_nodes:
        node_id = leveled["id"]
        node_level = leveled["level"]
        prereq_ids = prerequisites_by_id.get(node_id, [])

        # Gather unique doc_id/page source references
        seen_sources: set[str] = set()
        sources_out: list[dict] = []
        for src in leveled.get("sources", []):
            chunk_id = src.get("chunk_id", "")
            chunk = chunks_by_id.get(chunk_id)
            if chunk:
                key = f"{chunk.doc_id}:{chunk.page}"
                if key not in seen_sources:
                    seen_sources.add(key)
                    sources_out.append({"doc_id": chunk.doc_id, "page": chunk.page})

        # concept dict for difficulty
        concept_for_node = next((c for c in concept_dicts if c["id"] == node_id), {})
        node_difficulty = concept_for_node.get("difficulty", "intermediate")

        nodes_out.append(
            {
                "id": node_id,
                "title": leveled["name"],
                "concept_key": _slugify(leveled["name"]),
                "level": node_level,
                "difficulty": node_difficulty,
                "difficulty_badge": _DIFFICULTY_BADGE.get(node_difficulty, "⭐ Unknown"),
                "xp_reward": _xp_reward(node_level),
                "estimated_minutes": _estimated_minutes(node_level),
                "status": leveled["status"],
                "prerequisites": prereq_ids,
                "lesson": {
                    "summary": "",
                    "real_world_example": "",
                },
                "quiz": None,  # generated on-demand later
                "sources": sources_out,
                "position": {"x": 0, "y": 0},  # overwritten below
            }
        )

    # Calculate positions
    edges_out = [{"from": e["from"], "to": e["to"]} for e in validated["edges"]]
    nodes_out = calculate_node_positions(nodes_out, edges_out)

    tier0_ids = [n["id"] for n in nodes_out if n["level"] == 0]
    locked_ids = [n["id"] for n in nodes_out if n["level"] > 0]

    _emit("done", f"Pipeline complete — {len(nodes_out)} node(s), {len(edges_out)} edge(s).")

    return {
        "topic": f"Topic {topic_id[:8]}",  # will be overwritten from DB
        "topic_id": topic_id,
        "generated_at": datetime.utcnow().isoformat() + "Z",
        "document_count": len(file_paths),
        "nodes": nodes_out,
        "edges": edges_out,
        "dropped_edges": validated["dropped_edges"],
        "user_progress": {
            "completed_nodes": [],
            "unlocked_nodes": tier0_ids,
            "locked_nodes": locked_ids,
            "total_nodes": len(nodes_out),
            "percent_complete": 0,
        },
    }
