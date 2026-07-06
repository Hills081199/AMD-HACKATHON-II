"""One-time authoring tool: writes a small, real intro-ML document set into
data/source_docs/ for feat-009's pipeline runner (scripts/build_demo_dataset.py)
to ingest. These are genuine short documents (not filler/lorem-ipsum text) —
real, accurate explanations of the same intro-ML topics as the sample tree
in docs/concept-graph-pipeline.md, so the eventual real pipeline run has
plausible source material to extract concepts and infer prerequisites from.

Run once: `python scripts/generate_sample_docs.py`. The output files are
checked into data/source_docs/ so they don't need regenerating on every
machine that runs the pipeline.
"""

from __future__ import annotations

from pathlib import Path

_REPO_ROOT = Path(__file__).resolve().parents[1]
_OUTPUT_DIR = _REPO_ROOT / "data" / "source_docs"

_LINEAR_ALGEBRA_PAGES = [
    (
        "Vectors and Vector Spaces\n\n"
        "A vector is an ordered list of numbers representing a point or direction "
        "in space, e.g. [2, 5, -1]. In machine learning, a single data example is "
        "almost always represented as a vector: an image's pixels, a document's "
        "word counts, or a row of a spreadsheet. A vector space is the set of all "
        "vectors of a given dimension, closed under addition and scalar "
        "multiplication — the mathematical structure that lets us add two data "
        "points together or scale one by a constant."
    ),
    (
        "Matrices and Matrix Multiplication\n\n"
        "A matrix is a rectangular grid of numbers, and can be thought of as a "
        "batch of vectors stacked together — a full dataset of N examples with D "
        "features each is naturally an N-by-D matrix. Matrix multiplication "
        "combines two matrices by taking dot products of rows and columns; a "
        "neural network layer is, at its core, a matrix multiplication followed "
        "by a nonlinearity, applied to a batch of input vectors at once."
    ),
    (
        "The Dot Product\n\n"
        "The dot product of two vectors multiplies their corresponding entries "
        "and sums the results, producing a single number that measures how "
        "aligned the two vectors are. A large positive dot product means the "
        "vectors point in similar directions; a value near zero means they are "
        "roughly perpendicular. This single operation underlies cosine "
        "similarity, weighted sums in a neural network neuron, and projecting "
        "one vector onto another."
    ),
]

_CALCULUS_PAGES = [
    (
        "Derivatives and Rates of Change\n\n"
        "The derivative of a function measures how quickly its output changes as "
        "its input changes — the slope of the function at a point. For a loss "
        "function that scores how wrong a model's prediction is, the derivative "
        "with respect to a model parameter tells us whether increasing that "
        "parameter would make the loss better or worse, and by how much."
    ),
    (
        "Partial Derivatives and Gradients\n\n"
        "A model usually has many parameters at once, not just one. A partial "
        "derivative measures the rate of change with respect to a single "
        "parameter while holding all others fixed. Collecting every partial "
        "derivative into one vector gives the gradient, which points in the "
        "direction of steepest increase of the function — and its negative "
        "points toward the steepest decrease."
    ),
    (
        "The Chain Rule\n\n"
        "The chain rule lets us compute the derivative of a function that is "
        "built by composing simpler functions, by multiplying the derivatives of "
        "each step together. A deep neural network is exactly such a "
        "composition — layer after layer applied in sequence — so the chain rule "
        "is what makes it possible to compute how the final loss depends on a "
        "parameter buried many layers deep, which is what backpropagation does."
    ),
]

_PROBABILITY_SECTIONS = [
    ("Random Variables", (
        "A random variable assigns a number to the outcome of an uncertain "
        "event — the number of typos in an email, or whether an image contains "
        "a cat. Machine learning models routinely treat their own predictions as "
        "random variables, since the same input can plausibly map to several "
        "different outputs with different likelihoods."
    )),
    ("Probability Distributions", (
        "A probability distribution describes how likely each possible value of "
        "a random variable is. Discrete distributions assign a probability to "
        "each distinct outcome (like a coin flip); continuous distributions use "
        "a density function instead. A classifier's output layer typically "
        "produces a full probability distribution over class labels, not just a "
        "single guess."
    )),
    ("Expectation and Variance", (
        "The expectation of a random variable is its long-run average value, "
        "weighted by how likely each outcome is. Variance measures how spread "
        "out the values are around that average. Many training objectives — "
        "including the loss functions used to train neural networks — are "
        "themselves expectations, averaged over a batch of training examples."
    )),
]

_ML_SLIDES = [
    ("Loss Functions", (
        "A loss function scores how wrong a model's prediction is compared to "
        "the true answer.\nMean squared error penalizes large mistakes more "
        "than small ones — the difference is squared before averaging.\nThe "
        "training process searches for parameters that make this score as small "
        "as possible across the whole dataset."
    )),
    ("Gradient Descent", (
        "Gradient descent updates each parameter by taking a small step in the "
        "direction that most reduces the loss.\nThe learning rate controls how "
        "big that step is: too large overshoots the minimum, too small trains "
        "very slowly.\nRepeating this update thousands of times gradually walks "
        "the parameters toward a low-loss configuration."
    )),
    ("Model Evaluation Metrics", (
        "Accuracy alone can be misleading, especially when one class is much "
        "rarer than another.\nPrecision and recall separately measure how many "
        "predicted positives were correct, and how many actual positives were "
        "found.\nHeld-out validation data — examples the model never trained "
        "on — is what makes these metrics trustworthy."
    )),
    ("Controlling Overfitting", (
        "A model that memorizes its training data instead of learning general "
        "patterns will perform poorly on new examples — this is overfitting.\n"
        "Regularization techniques, such as penalizing large weights or randomly "
        "dropping units during training, discourage the model from relying too "
        "heavily on any single feature.\nComparing training loss against "
        "validation loss over time is the standard way to detect overfitting "
        "early."
    )),
]

_DEEP_LEARNING_SECTIONS = [
    ("Neural Network Architecture", (
        "A neural network is built from layers of simple units, each computing "
        "a weighted sum of its inputs. Stacking layers lets the network combine "
        "simple patterns learned in earlier layers into increasingly complex "
        "ones — edges into shapes, shapes into objects, for example, in an image "
        "classifier."
    )),
    ("Activation Functions", (
        "Without a nonlinearity between layers, stacking any number of matrix "
        "multiplications collapses into a single linear transformation, no more "
        "powerful than one layer. An activation function such as ReLU applies a "
        "simple nonlinear rule to each unit's output, which is what allows deep "
        "networks to approximate far more complex functions than a linear model "
        "ever could."
    )),
    ("Training via Backpropagation", (
        "Backpropagation computes the gradient of the loss with respect to "
        "every parameter in the network by applying the chain rule layer by "
        "layer, working backward from the output. Those gradients are then used "
        "by gradient descent to update every parameter at once, which is how a "
        "network with millions of parameters is trained efficiently."
    )),
]


def _write_pdf(path: Path, page_texts: list[str]) -> None:
    import fitz  # PyMuPDF

    document = fitz.open()
    for text in page_texts:
        page = document.new_page()
        page.insert_textbox(page.rect + (36, 36, -36, -36), text, fontsize=11)
    document.save(path)
    document.close()


def _write_docx(path: Path, sections: list[tuple[str, str]]) -> None:
    from docx import Document

    document = Document()
    for heading, body in sections:
        document.add_heading(heading, level=1)
        document.add_paragraph(body)
    document.save(path)


def _write_pptx(path: Path, slides: list[tuple[str, str]]) -> None:
    from pptx import Presentation

    presentation = Presentation()
    layout = presentation.slide_layouts[1]  # title + content
    for title, body in slides:
        slide = presentation.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text_frame.text = body
    presentation.save(path)


def main() -> None:
    _OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    _write_pdf(_OUTPUT_DIR / "linear_algebra_basics.pdf", _LINEAR_ALGEBRA_PAGES)
    _write_pdf(_OUTPUT_DIR / "calculus_for_ml.pdf", _CALCULUS_PAGES)
    _write_docx(_OUTPUT_DIR / "probability_for_ml.docx", _PROBABILITY_SECTIONS)
    _write_pptx(_OUTPUT_DIR / "ml_lecture_slides.pptx", _ML_SLIDES)
    _write_docx(_OUTPUT_DIR / "deep_learning_intro.docx", _DEEP_LEARNING_SECTIONS)
    print(f"Wrote 5 source documents to {_OUTPUT_DIR}")


if __name__ == "__main__":
    main()
